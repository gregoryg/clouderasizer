#Class for PPT plan. A PPT plan will auto generate a canned PPT based on a set of metrics that were pulled via a CollectionPlan
import os, metrics, collectionplan, logging
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.util import Pt
from collections import deque, Counter, OrderedDict

def create_impala_table_slide(prs,slide_title,data,headers):
    # change title and properties

    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = slide.shapes.title
    title.text = slide_title
    title.text_frame.paragraphs[0].font.name = "calibri"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    shapes = slide.shapes
    rows = len(data) + 1
    cols = len(headers)
    left = Inches(2.0)
    top = Inches(1.5)
    width = Inches(8.0)
    height = Inches(0.8)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    # set column widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(4.0)

    # write column headings
    for i in range(cols):
        table.cell(0,i).text = headers[i]
        
    # write body cells
    i = 1
    j = 0
    for row in data:
        for col in row:
            table.cell(i,j).text = str(col)
            table.cell(i,j).text_frame.paragraphs[0].font.name = "calibri"
            table.cell(i,j).text_frame.paragraphs[0].font.size = Pt(11)
            j+=1
        i+=1
        j=0
    return

def create_impala_query_slide(prs,chart_series_title, categories, values):
    # change title and properties
   
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = slide.shapes.title
    title.text = chart_series_title
    title.text_frame.paragraphs[0].font.name = "calibri"
    title.text_frame.paragraphs[0].font.size = Pt(28)

    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = categories
    chart_series_title = chart_series_title
    chart_data.add_series(chart_series_title, values)


    # add chart to slide --------------------
    x, y, cx, cy = Inches(0.5), Inches(1.4), Inches(9), Inches(6.0)
    chart =  slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    chart.chart_style = 10
    # sets chart plot to all be the same color (default in PPT is all different colors by category)
    plot = chart.chart.plots[0]
    plot.vary_by_categories = True

    # change axis properties
    category_axis = chart.chart.category_axis
    value_axis = chart.chart.value_axis

    # change axis font properties
    category_axis.tick_labels.font.name = "calibri"
    category_axis.tick_labels.font.size = Pt(11) 
    value_axis.tick_labels.font.name = "calibri"
    value_axis.tick_labels.font.size = Pt(11)

    # change legend properties
    chart.chart.has_legend = True
    chart.chart.legend.include_in_layout = False
    chart.chart.legend.font.name = "calibri"
    chart.chart.legend.font.size = Pt(11)
    return 

#Creates a series of slides with Impala query metrics
def create_impala_query_slides(prs,metric):
    timestamps = deque()
    values = deque()
    usernames = deque()
    statements = deque()
    databases = deque()
    values_statements = {}

    for query in metric['timeSeries']:
        data = query['data'][0]
        metadata = query['metadata']
        attributes = metadata['attributes']
        timestamp = data['timestamp']
        value = data['value']
        username = attributes['user']
        statement = attributes['statement']
        database = attributes['database']
        #append values to the lists in order
        values.append(value)
        usernames.append(username)
        statements.append(statement)
        databases.append(database)
        timestamps.append(timestamp)
        values_statements[value] = statement
    
    categories = deque()
    points = deque()

    #number of queries by user
    user_counter = Counter(usernames)
    user_counter_list = user_counter.items()
    for item in user_counter_list:
        categories.append(item[0])
        points.append(item[1])
    create_impala_query_slide(prs,"Impala Queries By User", categories, points)
    
    #number of queries by duration
    categories = deque(["0-9ms","10-99ms","100-999ms","1-59sec","1-9min","10-59min","1hr+"])
    points = deque()
    less_than_10ms = 0
    less_than_100ms = 0
    less_than_1s = 0
    less_than_1m = 0
    less_than_10m = 0
    less_than_1hr = 0
    over_1hr = 0
    query_duration_counter = Counter(values)
    query_duration_list = query_duration_counter.items()
    for item in query_duration_list:
        point = item[0]
        count = item[1]
        if point < 10:
            less_than_10ms += count
        elif point < 100:
            less_than_100ms += count
        elif point < 1000:
            less_than_1s += count
        elif point < 60000:
            less_than_1m += count
        elif point < 600000:
            less_than_10m += count
        else:
            over_1hr += count
    points = deque([less_than_10ms,less_than_100ms,less_than_1s,less_than_1m,less_than_10m,less_than_1hr,over_1hr])
    create_impala_query_slide(prs,"Impala Queries By Duration", categories, points)

    #top 10 queries by duration
    ordered_queries = OrderedDict(reversed(sorted(values_statements.items(), key=lambda x: x[0]))).items()
    top_10_queries = ordered_queries[0:10]
    headers = deque(["Duration(ms)","Query"])
    create_impala_table_slide(prs,"Top 10 Queries By Duration",top_10_queries,headers)
    return

#Create a metric slide
def create_metric_slide(prs,metric):
    #metric info
    category =  metric['timeSeries'][0]['metadata']['attributes']['category']
    #determine if service, and send to IMPALA_QUERY function if an IMPALA_QUERY slide
    if category == 'SERVICE':
        service_name = metric['timeSeries'][0]['metadata']['attributes']['serviceDisplayName']
    elif category == 'IMPALA_QUERY':
        create_impala_query_slides(prs,metric)
        return
    else:
        service_name = 'Cluster'

    metric_name = metric['timeSeries'][0]['metadata']['metricName']
    
    formatted_metric_name = ''
    #format the metric name and capitalize first letters
    metric_name_arr = metric_name.split("_")
    for word in metric_name_arr:
        formatted_metric_name += word.title() + " "
    
    metric_name = formatted_metric_name

    value_access_name = 'max'
         
    unit_type = metric['timeSeries'][0]['metadata']['unitNumerators'][0]
    numerator = metric['timeSeries'][0]['metadata']['unitNumerators'][0]
    denominators = metric['timeSeries'][0]['metadata']['unitDenominators']
    
    timestamps = deque()
    values = deque()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for point in metric['timeSeries'][0]['data']:
        formatted_timestamp = point['timestamp'].split("T")
        timestamps.append(formatted_timestamp[0])
        if unit_type == 'bytes':
            #value_in_mb = float(point['aggregateStatistics']['max'])/(1024*1024))
            value_in_gb = float(point['aggregateStatistics'][value_access_name])/(1024*1024*1024)
            values.append(value_in_gb)
            numerator = 'gigabytes' 
        else:
            values.append(point['aggregateStatistics'][value_access_name])
    
    #determine what our measurement (IE bytes/sec or just queries)
    if len(denominators) > 0:
        denominator = metric['timeSeries'][0]['metadata']['unitDenominators'][0]
        measurement = numerator + '/' + denominator
    else:
        denominator = 'None'
        measurement = numerator

    # change title and properties
    title = slide.shapes.title
    title.text = service_name + " " + metric_name
    title.text_frame.paragraphs[0].font.name = "calibri"
    title.text_frame.paragraphs[0].font.size = Pt(28)

    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = timestamps
    chart_series_title = measurement
    chart_data.add_series(chart_series_title, values)
    

    # add chart to slide --------------------
    x, y, cx, cy = Inches(0.5), Inches(1.4), Inches(9), Inches(6.0)
    chart =  slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data)
    
    # sets chart plot to all be the same color (default in PPT is all different colors by category)
    plot = chart.chart.plots[0]
    plot.vary_by_categories = False

    # change axis properties
    category_axis = chart.chart.category_axis
    value_axis = chart.chart.value_axis

    # change axis font properties
    category_axis.tick_labels.font.name = "calibri"
    category_axis.tick_labels.font.size = Pt(11)
    value_axis.tick_labels.font.name = "calibri"
    value_axis.tick_labels.font.size = Pt(11)
   
    # change legend properties
    chart.chart.has_legend = True
    chart.chart.legend.include_in_layout = False
    chart.chart.legend.font.name = "calibri"
    chart.chart.legend.font.size = Pt(11)
    return

#Create a service slide
def create_service_slide():
    return

#Create a summary slide
def create_summary_slide():
    return

#Create a title slide
def create_title_slide(prs):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Clouderasizer"
    subtitle.text = "Created by Brandon Kvarda"
    return 

#Parse pptplan
def parse_pptplan(pptplan):
    return

#Create a presentation
def create_ppt(collection_zip,output_dir):
    #Unzip the collection
    collection_dir = collectionplan.unzip_collection(collection_zip,output_dir) 
    #Parse the collection
    collection = deque()
    
    for file in os.listdir(collection_dir):
        full_file_path = collection_dir + '/' + file
        collection.append(metrics.read_from_json(full_file_path))
    
    #Instantiate presentation
    prs = Presentation()
    logging.info("Creating Presentation")    
    #create the title slide
    create_title_slide(prs)
    #create metric slide for each metric
    for metric in collection:
       #structure is different depending on whether it was pulled using CM API library for python or directly from CM API (IE using CURL)
        if 'items' in metric:
            metric = metric['items'][0]

        metric_name = metric['timeSeries'][0]['metadata']['metricName']
        logging.info("Creating Slide For: " + metric_name)
        create_metric_slide(prs,metric)    
    
    #save output
    output_file = output_dir + '/' + 'test.pptx'
    logging.info("Saving presentation at location: " + output_file)
    prs.save(output_file)


